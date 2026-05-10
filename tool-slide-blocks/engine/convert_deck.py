# -*- coding: utf-8 -*-
"""
色系风格转换工具
将整份 PPT 从深色底转为浅色底，或从浅色底转为深色底。

核心原理：与组装流程完全相同——逐页通过模板引擎重建，
模板提供背景框架，源内容以保留源格式方式粘贴，再修复文字颜色。

由 Claude Code 根据自然语言对话生成调用脚本，例如：
  from convert_deck import convert
  convert("D:/输入.pptx", to="light")   # 深色底 → 浅色底，自动选模板
  convert("D:/输入.pptx", to="dark")    # 浅色底 → 深色底，自动选模板
  convert("D:/输入.pptx", to="dark",
          template_path="D:/Claude/SlideBlocks/模板/科技风（深色底）.pptx")
"""

from pathlib import Path
import sys

_ENGINE = Path(__file__).parent
_ROOT   = _ENGINE.parent  # SlideBlocks 根目录

import engine.assemble_template as _at

TEMPLATE_DIR = _ROOT / "模板"


# ─── 模板选择 ─────────────────────────────────────────────────────────────────

def _pick_template(to):
    """从 模板/ 文件夹自动选匹配方向的第一个模板。"""
    keyword = "深色底" if to == "dark" else "浅色底"
    matches = [p for p in TEMPLATE_DIR.iterdir()
               if p.suffix.lower() in ('.pptx', '.potx')
               and keyword in p.name
               and not p.name.startswith('~')]
    if not matches:
        raise FileNotFoundError(
            f"模板文件夹中没有找到含 [{keyword}] 的模板：{TEMPLATE_DIR}"
        )
    if len(matches) > 1:
        print(f"  [模板] 找到多个 {keyword} 模板：{[p.name for p in matches]}，使用：{matches[0].name}")
    return matches[0]


# ─── 页面类型检测（用 python-pptx 静态分析源文件） ────────────────────────────

def _extract_title(slide):
    """从幻灯片中提取最可能的标题文字（位置靠上、文字不过长）。"""
    candidates = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text or len(text) > 80:
            continue
        candidates.append((shape.top, len(text), text))
    if not candidates:
        return ""
    candidates.sort()           # 按 top 位置升序，取最靠上的
    return candidates[0][2][:60]


def _is_transition_slide(slide):
    """
    判断是否为过渡/章节标题页：
    无图片、无图表、文字形状 ≤ 3 个、总文字 ≤ 60 字。
    """
    for shape in slide.shapes:
        if getattr(shape, 'shape_type', 0) == 13:   # PICTURE
            return False
        try:
            if shape.has_chart:
                return False
        except Exception:
            pass
    text_shapes = [s for s in slide.shapes
                   if s.has_text_frame and s.text_frame.text.strip()]
    all_text = ' '.join(s.text_frame.text.strip() for s in text_shapes)
    return len(text_shapes) <= 3 and len(all_text) <= 60


def _is_end_slide(slide):
    """判断是否为封底（文字极少，通常只有 logo/标语）。"""
    text_shapes = [s for s in slide.shapes
                   if s.has_text_frame and s.text_frame.text.strip()]
    all_text = ' '.join(s.text_frame.text.strip() for s in text_shapes)
    return len(all_text) <= 30


def _is_cover_slide(slide):
    """
    判断是否为封面页：无图片/图表/表格，文字形状 ≤ 3，总文字 ≤ 60 字。
    标准与 _is_transition_slide 相同，封面和过渡页结构上都是"标题+副标题"式极简布局。
    如果第一页内容更丰富（有详细正文），应当作普通内容页色系转换，而不是替换为模板封面。
    """
    for shape in slide.shapes:
        if getattr(shape, 'shape_type', 0) == 13:   # PICTURE
            return False
        try:
            if shape.has_chart:
                return False
        except Exception:
            pass
        try:
            if shape.has_table:
                return False
        except Exception:
            pass
    text_shapes = [s for s in slide.shapes
                   if s.has_text_frame and s.text_frame.text.strip()]
    all_text = ' '.join(s.text_frame.text.strip() for s in text_shapes)
    return len(text_shapes) <= 3 and len(all_text) <= 60


# ─── PLAN 自动生成 ────────────────────────────────────────────────────────────

def _auto_plan(src_path, to):
    """
    扫描源文件，自动生成 PLAN：
    - 第 1 页 → 模板封面（P1），提取标题
    - 最后 1 页（若内容极少）→ 模板封底（P5）
    - 中间过渡页（少文字、无图表）→ 模板过渡页（P2），提取章节标题
    - 其余 → 内容页（src + page，由引擎自动选 P3/P4）

    to 参数用于判断颜色修复方向：
    - 文件名含方向关键词时引擎会自动触发，无需额外字段
    - 文件名不含关键词时，根据 to 显式注入 fix_colors / fix_colors_dark，
      确保用户明确说"转成浅色/深色"时颜色修复一定生效

    关于组装规则（开篇不加过渡页 / 后续章节加过渡页）：
    源文件里的过渡页已经在原来的位置，直接映射为模板 P2 即可，
    无需 assemble_template 里的"开篇不加过渡页"规则干预。
    """
    from pptx import Presentation
    src_str = str(Path(src_path).resolve())
    src_name = Path(src_path).name

    # 文件名没有方向关键词时，根据 to 显式注入颜色修复标记
    needs_light_fix = to == "light" and "深色底" not in src_name
    needs_dark_fix  = to == "dark"  and "浅色底" not in src_name

    prs = Presentation(src_str)
    n   = len(prs.slides)
    plan = []

    for i, slide in enumerate(prs.slides):
        pn = i + 1

        if i == 0 and _is_cover_slide(slide):
            # 第一页是极简封面（无图表/表格，文字 ≤ 3 个且总字数 ≤ 60）→ 替换为模板封面
            title = _extract_title(slide)
            plan.append({"template_page": 1, "replace_title": title})

        elif i == 0:
            # 第一页有实质内容 → 当普通内容页处理（保留原始布局，只做色系转换）
            item = {"src": src_str, "page": pn}
            if needs_light_fix:
                item["fix_colors"] = True
            elif needs_dark_fix:
                item["fix_colors_dark"] = True
            plan.append(item)

        elif i == n - 1 and _is_end_slide(slide):
            # 最后一页且内容极少 → 封底
            plan.append({"template_page": 5})

        elif _is_transition_slide(slide):
            # 过渡/章节标题页 → 模板过渡页
            title = _extract_title(slide)
            plan.append({"template_page": 2, "replace_title": title})

        else:
            # 普通内容页 → 套模板背景 + 保留源内容（引擎自动选 P3/P4）
            item = {"src": src_str, "page": pn}
            if needs_light_fix:
                item["fix_colors"] = True
            elif needs_dark_fix:
                item["fix_colors_dark"] = True
            plan.append(item)

    return plan


# ─── 主入口 ───────────────────────────────────────────────────────────────────

def convert(src, output_name=None, to='light', template_path=None):
    """
    色系转换主函数。

    参数：
      src           源 PPTX 路径（字符串或 Path）
      output_name   输出文件名（不含扩展名，默认：源文件名 + _浅色底 / _深色底）
      to            'light'  深色底 → 浅色底
                    'dark'   浅色底 → 深色底
      template_path 指定模板路径（默认自动从 模板/ 文件夹选匹配方向的第一个）

    返回：输出文件 Path。
    """
    src = Path(src)
    if not src.exists():
        raise FileNotFoundError(f"源文件不存在：{src}")
    if to not in ('light', 'dark'):
        raise ValueError(f"to 只能是 'light' 或 'dark'，收到：{to!r}")

    # 选模板
    if template_path is None:
        template_path = _pick_template(to)
    template_path = Path(template_path)

    # 输出文件名
    if output_name is None:
        suffix = "_浅色底" if to == "light" else "_深色底"
        output_name = src.stem + suffix

    direction_label = "深色底 → 浅色底" if to == "light" else "浅色底 → 深色底"
    print(f"[色系转换] {src.name}")
    print(f"  方向：{direction_label}")
    print(f"  模板：{Path(template_path).name}")

    # 自动生成 PLAN
    plan = _auto_plan(src, to)
    n_content    = sum(1 for p in plan if "src" in p)
    n_transition = sum(1 for p in plan if p.get("template_page") == 2)
    print(f"  页数：共 {len(plan)} 页（内容页 {n_content}，过渡页 {n_transition}，封面1，封底1）\n")

    # 构建完整输出路径（与源文件同目录）
    output_full_path = str(src.parent / (output_name + ".pptx"))

    # 调用组装引擎（颜色修复在引擎内自动触发）
    _at.assemble(plan, output_full_path, str(template_path))
    return Path(output_full_path)
