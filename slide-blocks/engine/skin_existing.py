# -*- coding: utf-8 -*-
"""
skin_existing.py — 给已有手工拼装 PPT 统一换皮

用法：
    python skin_existing.py <输入文件> [选项]

选项：
    --template   模板路径（默认：模板/深色底模板.pptx）
    --keep       保留不换皮的页码，逗号分隔（默认：首页和末页）
                 例：--keep 1,10
    --output     输出文件名（不含扩展名，默认：原文件名_skinned）

示例：
    python skin_existing.py 测试.PPTX
    python skin_existing.py 测试.PPTX --keep 1,10 --template 模板/浅色底模板.pptx
"""

import sys
import argparse
from pathlib import Path
from pptx import Presentation

# 引入核心组装引擎
sys.path.insert(0, str(Path(__file__).parent))  # engine/ 内互相引用
from assemble_template import assemble, TEMPLATE_PATH


def build_plan(input_path: Path, keep_pages: list[int], template_path: Path) -> list[dict]:
    """根据输入 PPT 和保留页列表，生成 assemble() 所需的 plan。"""
    prs = Presentation(str(input_path))
    total = len(prs.slides)
    print(f"[扫描] {input_path.name}  共 {total} 页")
    print(f"[保留] 第 {keep_pages} 页（整页复制，不换皮）")
    print(f"[模板] {template_path.name}\n")

    plan = []
    for page_num in range(1, total + 1):
        if page_num in keep_pages:
            # 整页保留：直接复制源文件该页，不套模板
            plan.append({
                "copy_slide": str(input_path),
                "page": page_num,
            })
        else:
            # 内容页：套模板框架（深色/浅色背景 + 标题栏），保留内容形状颜色
            plan.append({
                "src": str(input_path),
                "page": page_num,
            })

    return plan


def main():
    parser = argparse.ArgumentParser(description="给已有 PPT 统一换皮")
    parser.add_argument("input", help="输入 PPT 文件路径")
    parser.add_argument("--template", default="", help="模板文件路径（留空则交互选择）")
    parser.add_argument("--keep", default="", help="保留页码（逗号分隔，默认首页和末页）")
    parser.add_argument("--output", default="", help="输出文件名（不含扩展名）")
    args = parser.parse_args()

    input_path = Path(args.input).resolve()
    if not input_path.exists():
        print(f"[错误] 文件不存在: {input_path}")
        sys.exit(1)

    # 解析保留页
    prs = Presentation(str(input_path))
    total = len(prs.slides)
    if args.keep:
        keep_pages = [int(p.strip()) for p in args.keep.split(",") if p.strip()]
    else:
        keep_pages = [1, total]  # 默认保留首页和末页

    # 模板路径：未指定则扫描模板文件夹让用户选择
    template_dir = Path(__file__).parent.parent / "模板"
    if args.template:
        template_path = Path(args.template).resolve()
        if not template_path.exists():
            print(f"[错误] 模板不存在: {template_path}")
            sys.exit(1)
    else:
        templates = sorted(template_dir.glob("*.pptx"))
        if not templates:
            print(f"[错误] 模板文件夹为空：{template_dir}")
            print("请在模板文件夹中放入 .pptx 模板文件（须包含5页：封面、过渡页、带标题内容页、无标题内容页、封底）")
            sys.exit(1)
        print("可用模板：")
        for idx, t in enumerate(templates, 1):
            print(f"  {idx}. {t.stem}")
        if len(templates) == 1:
            template_path = templates[0]
            print(f"→ 自动选择唯一模板：{template_path.stem}\n")
        else:
            choice = input(f"请选择模板编号（1-{len(templates)}）：").strip()
            try:
                template_path = templates[int(choice) - 1]
            except (ValueError, IndexError):
                print("[错误] 无效选择")
                sys.exit(1)
            print()

    # 临时覆盖引擎的 TEMPLATE_PATH
    import assemble_template
    assemble_template.TEMPLATE_PATH = template_path

    # 输出文件名
    output_name = args.output or (input_path.stem + "_skinned")

    # 生成 plan 并执行
    plan = build_plan(input_path, keep_pages, template_path)
    assemble(plan, output_name)


if __name__ == "__main__":
    main()
