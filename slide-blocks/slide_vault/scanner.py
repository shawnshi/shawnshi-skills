"""
scanner.py - PPT 扫描 & 文字提取模块 (Multiprocessing & Batch Insert Optimized)
"""

import os
import sys
import hashlib
import sqlite3
import json
from datetime import datetime
from pathlib import Path
from concurrent.futures import ProcessPoolExecutor, as_completed

# Windows 终端强制 UTF-8 输出
if sys.stdout.encoding != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .config import get_db_path, get_materials_dir, load_config
from .vector_store import init_vector_db
from engine.thumbnail_exporter import export_thumbnails


# ─── 数据库初始化 ────────────────────────────────────────────────

def init_db():
    conn = sqlite3.connect(get_db_path())
    conn.execute("""
        CREATE TABLE IF NOT EXISTS slides (
            id            INTEGER PRIMARY KEY AUTOINCREMENT,
            file_path     TEXT NOT NULL,
            file_name     TEXT NOT NULL,
            file_hash     TEXT NOT NULL,
            slide_index   INTEGER NOT NULL,
            title         TEXT,
            body_text     TEXT,
            shape_count   INTEGER,
            has_image     BOOLEAN,
            has_chart     BOOLEAN,
            file_mtime    TEXT,
            indexed_at    TEXT,
            thumbnail_path TEXT,
            UNIQUE(file_hash, slide_index)
        )
    """)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS tags (
            slide_id      INTEGER PRIMARY KEY,
            scene         TEXT,
            content_type  TEXT,
            industries    TEXT,
            keywords      TEXT,
            quality_score INTEGER,
            summary       TEXT,
            tagged_at     TEXT,
            layout_type   TEXT,
            FOREIGN KEY(slide_id) REFERENCES slides(id)
        )
    """)
    conn.commit()
    conn.close()
    init_vector_db(get_db_path())
    print(f"[DB] 数据库及向量表已初始化：{get_db_path()}")


# ─── 文件 Hash ────────────────────────────────────────────────────

def file_hash(path: Path) -> str:
    h = hashlib.md5()
    with open(path, "rb") as f:
        while chunk := f.read(8192):
            h.update(chunk)
    return h.hexdigest()


# ─── 单页内容提取 ────────────────────────────────────────────────

def extract_slide_content(slide) -> dict:
    texts = []
    title = None
    has_image = False
    has_chart = False
    shape_count = len(slide.shapes)

    for shape in slide.shapes:
        # 标题优先取 placeholder 类型为 title 的
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                try:
                    ph = shape.placeholder_format
                    if ph and ph.idx == 0:  # idx=0 通常是标题
                        title = text
                    else:
                        texts.append(text)
                except Exception:
                    texts.append(text)
        elif hasattr(shape, "image"):
            has_image = True

        if shape.shape_type == 13:  # PICTURE
            has_image = True
        elif shape.shape_type == 19: # TABLE
            has_chart = True

        try:
            if shape.has_chart:
                has_chart = True
        except Exception:
            pass

    # 如果 title 没有从 placeholder 里找到，退而用第一段文字
    if not title and texts:
        title = texts[0]
        texts = texts[1:]

    body_text = "\n".join(texts).replace("\xa0", " ")
    return {
        "title": title,
        "body_text": body_text,
        "shape_count": shape_count,
        "has_image": has_image,
        "has_chart": has_chart,
    }


# ─── 独立 Worker 扫描单文件 ────────────────────────────────────────

def scan_file_worker(args: tuple) -> list[tuple]:
    """独立进程中解析 PPT 文件。返回供批量入库的元组列表。"""
    pptx_path_str, fname, fhash, fmtime, expected_count = args
    try:
        prs = Presentation(pptx_path_str)
        total = len(prs.slides)
    except Exception as e:
        print(f"  [!] 无法打开 {fname}：{e}")
        return []

    if expected_count is not None and expected_count >= total:
        # 已完整索引
        return []

    results = []
    now = datetime.now().isoformat()

    for i, slide in enumerate(prs.slides, start=1):
        try:
            content = extract_slide_content(slide)
            results.append((
                pptx_path_str,
                fname,
                fhash,
                i,
                content["title"],
                content["body_text"],
                content["shape_count"],
                content["has_image"],
                content["has_chart"],
                fmtime,
                now
            ))
        except Exception as e:
            print(f"  [!] {fname} 第{i}页提取失败：{e}")

    return results


# ─── 批量扫描目录 ────────────────────────────────────────────────

def scan_directory(directory: Path = None):
    if directory is None:
        directory = get_materials_dir()

    cfg = load_config()
    exclude_dirs = set(cfg.get("exclude_dirs", []))

    # 预览图根目录设置为素材库同级的 .thumbnails
    thumb_root = Path(directory).parent / ".thumbnails"
    thumb_root.mkdir(parents=True, exist_ok=True)

    init_db()
    conn = sqlite3.connect(get_db_path())
    
    # 提前查询已有记录的页数
    existing_counts = dict(
        conn.execute("SELECT file_hash, COUNT(*) FROM slides GROUP BY file_hash").fetchall()
    )

    # 收集文件，跳过排除目录
    pptx_files = []
    for f in directory.rglob("*.pptx"):
        if any(part in exclude_dirs for part in f.parts):
            continue
        # 排除以 ~$ 开头的临时文件
        if f.name.startswith("~"):
            continue
        pptx_files.append(f)

    print(f"\n[扫描] 发现 {len(pptx_files)} 个文件（已排除：{exclude_dirs}）\n")

    tasks = []
    for pptx_path in pptx_files:
        fhash = file_hash(pptx_path)
        fmtime = datetime.fromtimestamp(pptx_path.stat().st_mtime).isoformat()
        expected_count = existing_counts.get(fhash)
        tasks.append((str(pptx_path), pptx_path.name, fhash, fmtime, expected_count))

    total_slides = 0
    # 降级为单进程顺序扫描，确保 COM (win32com) 环境稳定性
    for task in tasks:
        pptx_path_str, fname, fhash, fmtime, expected_count = task
        try:
            # Step 1: 提取文本
            results = scan_file_worker(task)
            if not results:
                if expected_count is not None and expected_count > 0:
                    print(f"  [跳过] {fname}（已索引 {expected_count} 页）")
                continue
            
            print(f"[处理] {fname} → 提取 {len(results)} 页文本")

            # Step 2: 导出预览图
            pptx_p = Path(pptx_path_str)
            thumb_folder = export_thumbnails(pptx_p, thumb_root)
            
            # Step 3: 批量入库
            conn.execute("BEGIN TRANSACTION")
            final_rows = []
            for row in results:
                s_idx = row[3]
                t_path = thumb_folder / f"slide_{s_idx:03d}.jpg" if thumb_folder else None
                final_rows.append(row + (str(t_path) if t_path else None,))

            conn.executemany(
                """
                INSERT OR IGNORE INTO slides
                (file_path, file_name, file_hash, slide_index,
                 title, body_text, shape_count, has_image, has_chart,
                 file_mtime, indexed_at, thumbnail_path)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                """,
                final_rows
            )
            conn.commit()
            total_slides += len(results)
                
        except Exception as e:
            print(f"  [!] 解析 {fname} 异常：{e}")

    conn.close()
    print(f"\n[完成] 共新增 {total_slides} 条带视觉索引的幻灯片记录")


# ─── 预览提取结果 ────────────────────────────────────────────────

def preview_results(limit: int = 10):
    conn = sqlite3.connect(get_db_path())
    rows = conn.execute(
        """
        SELECT slide_index, title, body_text, shape_count, has_image, has_chart
        FROM slides
        ORDER BY slide_index
        LIMIT ?
        """,
        (limit,),
    ).fetchall()
    conn.close()

    print(f"\n{'='*60}")
    print(f"前 {limit} 页提取结果预览")
    print(f"{'='*60}")
    for row in rows:
        idx, title, body, shapes, img, chart = row
        body_preview = (body or "")[:100].replace("\n", " | ")
        print(f"\n【第 {idx} 页】")
        print(f"  标题：{title or '（无）'}")
        print(f"  正文：{body_preview or '（无文字）'}")
        print(f"  元素数：{shapes}  含图片：{'是' if img else '否'}  含图表：{'是' if chart else '否'}")


if __name__ == "__main__":
    scan_directory()
    preview_results(limit=20)
