"""
<!-- Standard Header -->
@Input: outline.md (Narrative Blueprint)
@Output: deck.pptx
@Phase: Phase 4 - Narrative Rendering & Assembly
@Maintenance Protocol: Refactor plan V9.0 - Narrative Engine.
"""
import os
import re
import sys
import argparse
from pptx import Presentation
from layout_engine import LayoutEngine

def parse_narrative_blueprint(content):
    # 1. Extract Style Instructions
    style_match = re.search(r'<STYLE_INSTRUCTIONS>([\s\S]*?)</STYLE_INSTRUCTIONS>', content)
    style_info = style_match.group(1).strip() if style_match else ""

    # 2. Extract Slides
    slides_data = []
    slides_raw = re.split(r'Page \d+:', content)
    if slides_raw and not slides_raw[0].strip():
        slides_raw = slides_raw[1:]
        
    for slide_raw in slides_raw:
        data = {}
        # Narrative extraction
        goal = re.search(r'// NARRATIVE GOAL\s*(.*?)\n', slide_raw)
        headline = re.search(r'Headline:\s*(.*?)\n', slide_raw)
        sub_headline = re.search(r'Sub-headline:\s*(.*?)\n', slide_raw)
        body = re.search(r'Body/Data:([\s\S]*?)// VISUAL', slide_raw)
        visual = re.search(r'// VISUAL\s*(.*?)\n', slide_raw)
        layout = re.search(r'// LAYOUT\s*(.*?)\n', slide_raw)
        script = re.search(r'// Script:([\s\S]*?)$', slide_raw)
        
        # Mapping to LayoutEngine compatible dictionary
        data['kicker'] = headline.group(1).strip() if headline else "Untitled Slide"
        data['lead_in'] = sub_headline.group(1).strip() if sub_headline else ""
        data['body'] = body.group(1).strip() if body else ""
        data['narrative_goal'] = goal.group(1).strip() if goal else ""
        data['visual_desc'] = visual.group(1).strip() if visual else ""
        data['layout_desc'] = layout.group(1).strip() if layout else ""
        data['script'] = script.group(1).strip() if script else ""
        
        # Add a placeholder for evidence/trust to keep LayoutEngine compatible
        data['evidence'] = ""
        data['trust_anchor'] = ""
        
        if data:
            slides_data.append(data)
            
    return style_info, slides_data

def main():
    parser = argparse.ArgumentParser(description="Mentat Narrative PPT Forger")
    parser.add_argument("dir", help="Directory containing outline.md")
    parser.add_argument("--output", "-o", help="Output filename")
    args = parser.parse_args()
    
    deck_dir = args.dir
    outline_file = os.path.join(deck_dir, "outline.md")
    
    # 1. Narrative Logic Audit
    SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
    validator_path = os.path.join(SCRIPT_DIR, "validator.py")
    
    print(f"--- Step 1: Mentat Narrative Audit ---")
    import subprocess
    audit_cmd = [sys.executable, validator_path, outline_file]
    result = subprocess.run(audit_cmd, capture_output=True, text=True)
    
    if result.returncode != 0:
        print(result.stdout)
        print("\n❌ ABORTING: Narrative audit failed.")
        sys.exit(1)
    print("✅ Narrative Audit Passed.")

    # 2. Parse Narrative Blueprint
    with open(outline_file, 'r', encoding='utf-8') as f:
        content = f.read()
    style_info, slides_data = parse_narrative_blueprint(content)
    
    # 3. Native Rendering
    print(f"\n--- Step 2: Native Object Rendering ({len(slides_data)} slides) ---")
    prs = Presentation()
    
    # Set 16:9 ratio (Widescreen)
    from pptx.util import Inches
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    engine = LayoutEngine(prs)
    
    for slide_data in slides_data:
        # Render and add notes (Narrative Goal + Script)
        engine.render_slide(slide_data)
        current_slide = prs.slides[-1]
        notes_slide = current_slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        
        full_notes = f"GOAL: {slide_data['narrative_goal']}\n\nVISUAL: {slide_data['visual_desc']}\n\nLAYOUT: {slide_data['layout_desc']}\n\nSCRIPT:\n{slide_data['script']}"
        notes_text_frame.text = full_notes
        
    # 4. Save
    output_name = args.output if args.output else f"{os.path.basename(deck_dir)}.pptx"
    output_path = os.path.join(deck_dir, output_name)
    prs.save(output_path)
    
    print(f"\n✅ NARRATIVE FORGING COMPLETE: {output_path}")

if __name__ == "__main__":
    main()
