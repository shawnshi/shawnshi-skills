#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
@Input:  Document directory, parallelism settings
@Output: extracted_content_part1.json (Text/ID mapping)
@Pos:    Kernel Layer. Responsible for high-fidelity text extraction across formats.

!!! Maintenance Protocol: If new file formats (e.g. .md, .txt) are needed, update SUPPORTED_EXTENSIONS.
!!! For medical OCR, consider integrating pytesseract in future iterations.

Document Summarizer - 文本提取脚本
支持 PDF、DOCX、PPTX、XLSX 文档的并发文本提取
"""
import os
import sys
import json
import hashlib
import argparse
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import Dict, List, Optional, Tuple

try:
    from pypdf import PdfReader
    from docx import Document
    from pptx import Presentation
    from openpyxl import load_workbook
    from tqdm import tqdm
except ImportError as e:
    print(f"❌ 缺少依赖包: {e}")
    print("💡 请运行: pip install -r requirements.txt")
    sys.exit(1)


class DocumentExtractor:
    """文档文本提取器"""

    SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.pptx', '.xlsx', '.xlsm'}

    def __init__(self, force: bool = False):
        self.force = force
        self.stats = {
            'total': 0,
            'extracted': 0,
            'skipped': 0,
            'failed': 0,
            'errors': []
        }

    @staticmethod
    def generate_file_id(file_path: str) -> str:
        """生成文件唯一ID（基于路径MD5）"""
        return hashlib.md5(file_path.encode('utf-8')).hexdigest()[:12]

    @staticmethod
    def has_metadata(file_path: Path) -> bool:
        """检查文件是否已有元数据"""
        if not file_path.exists():
            return False

        ext = file_path.suffix.lower()

        try:
            if ext == '.pdf':
                reader = PdfReader(str(file_path))
                metadata = reader.metadata
                return metadata and (metadata.get('/Subject') or metadata.get('/Keywords'))

            elif ext == '.docx':
                doc = Document(str(file_path))
                props = doc.core_properties
                return bool(props.subject or props.keywords)

            elif ext == '.pptx':
                prs = Presentation(str(file_path))
                props = prs.core_properties
                return bool(props.subject or props.keywords)

            elif ext in {'.xlsx', '.xlsm'}:
                wb = load_workbook(str(file_path), read_only=True)
                props = wb.properties
                wb.close()
                return bool(props.subject or props.keywords)

        except Exception:
            pass

        return False

    def extract_pdf(self, file_path: Path) -> Optional[str]:
        """提取PDF文本"""
        try:
            reader = PdfReader(str(file_path))
            text_parts = []

            # 提取前20页或全部页面
            max_pages = min(20, len(reader.pages))
            for page in reader.pages[:max_pages]:
                text = page.extract_text()
                if text:
                    text_parts.append(text)

            content = '\n'.join(text_parts)

            # 限制长度（最多10000字符）
            if len(content) > 10000:
                content = content[:10000] + '...'

            return content

        except Exception as e:
            raise Exception(f"PDF提取失败: {e}")

    def extract_docx(self, file_path: Path) -> Optional[str]:
        """提取Word文档文本"""
        try:
            doc = Document(str(file_path))
            text_parts = []

            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # 提取表格内容
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text for cell in row.cells)
                    if row_text.strip():
                        text_parts.append(row_text)

            content = '\n'.join(text_parts)

            # 限制长度
            if len(content) > 10000:
                content = content[:10000] + '...'

            return content

        except Exception as e:
            raise Exception(f"DOCX提取失败: {e}")

    def extract_pptx(self, file_path: Path) -> Optional[str]:
        """提取PowerPoint文本"""
        try:
            prs = Presentation(str(file_path))
            text_parts = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text_parts.append(shape.text)

            content = '\n'.join(text_parts)

            # 限制长度
            if len(content) > 10000:
                content = content[:10000] + '...'

            return content

        except Exception as e:
            raise Exception(f"PPTX提取失败: {e}")

    def extract_xlsx(self, file_path: Path) -> Optional[str]:
        """提取Excel文本"""
        try:
            wb = load_workbook(str(file_path), read_only=True, data_only=True)
            text_parts = []

            # 最多处理前5个工作表
            for sheet_name in list(wb.sheetnames)[:5]:
                sheet = wb[sheet_name]
                text_parts.append(f"[工作表: {sheet_name}]")

                # 最多读取前100行
                row_count = 0
                for row in sheet.iter_rows(values_only=True):
                    if row_count >= 100:
                        break

                    row_text = ' | '.join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        text_parts.append(row_text)
                        row_count += 1

            wb.close()
            content = '\n'.join(text_parts)

            # 限制长度
            if len(content) > 10000:
                content = content[:10000] + '...'

            return content

        except Exception as e:
            raise Exception(f"XLSX提取失败: {e}")

    def extract_file(self, file_path: Path) -> Optional[Dict]:
        """提取单个文件的文本内容"""
        ext = file_path.suffix.lower()

        # 检查是否已有元数据（非强制模式）
        if not self.force and self.has_metadata(file_path):
            self.stats['skipped'] += 1
            return None

        try:
            # 根据文件类型选择提取方法
            if ext == '.pdf':
                content = self.extract_pdf(file_path)
            elif ext == '.docx':
                content = self.extract_docx(file_path)
            elif ext == '.pptx':
                content = self.extract_pptx(file_path)
            elif ext in {'.xlsx', '.xlsm'}:
                content = self.extract_xlsx(file_path)
            else:
                return None

            if not content or not content.strip():
                raise Exception("提取内容为空")

            file_id = self.generate_file_id(str(file_path))
            self.stats['extracted'] += 1

            return {
                'id': file_id,
                'filename': str(file_path),
                'content': content.strip()
            }

        except Exception as e:
            self.stats['failed'] += 1
            self.stats['errors'].append({
                'file': str(file_path),
                'error': str(e)
            })
            return None

    def find_documents(self, directory: Path) -> List[Path]:
        """递归查找所有支持的文档"""
        documents = []

        for ext in self.SUPPORTED_EXTENSIONS:
            for file_path in directory.rglob(f'*{ext}'):
                # Ignore hidden and temporary files
                if not file_path.name.startswith('.') and not file_path.name.startswith('~$'):
                    documents.append(file_path)

        return sorted(documents)

    def extract_batch(self, directory: Path, workers: int = 5) -> Tuple[List[Dict], Dict]:
        """批量提取文档"""
        documents = self.find_documents(directory)
        self.stats['total'] = len(documents)

        if not documents:
            print(f"❌ 在目录 {directory} 中未找到支持的文档")
            return [], {}

        print(f"📁 找到 {len(documents)} 个文档文件")
        print(f"🔧 支持的格式: {', '.join(self.SUPPORTED_EXTENSIONS)}")

        if not self.force:
            print(f"⏭️  将跳过已有元数据的文件（使用 --force 强制重新提取）")

        extracted_data = []
        file_id_mapping = {}

        # 并发提取
        with ThreadPoolExecutor(max_workers=workers) as executor:
            futures = {executor.submit(self.extract_file, doc): doc for doc in documents}

            with tqdm(total=len(documents), desc="提取进度", unit="文件") as pbar:
                for future in as_completed(futures):
                    result = future.result()
                    if result:
                        extracted_data.append(result)
                        file_id_mapping[result['id']] = result['filename']
                    pbar.update(1)

        return extracted_data, file_id_mapping


def main():
    parser = argparse.ArgumentParser(
        description='文档文本提取工具 - 支持 PDF/DOCX/PPTX/XLSX',
        formatter_class=argparse.RawDescriptionHelpFormatter
    )

    parser.add_argument('--dir', required=True, help='文档目录路径')
    parser.add_argument('--workers', type=int, default=5, help='并行线程数（默认: 5）')
    parser.add_argument('--force', action='store_true', help='强制重新提取所有文件')
    parser.add_argument('--output', default='extracted_content_part1.json', help='输出文件名')
    parser.add_argument('--mapping', default='file_id_mapping.json', help='文件映射文件名')

    args = parser.parse_args()

    # 验证目录
    directory = Path(args.dir)
    if not directory.exists() or not directory.is_dir():
        print(f"❌ 目录不存在: {directory}")
        return 1

    print("\n" + "="*60)
    print("📄 文档文本提取工具")
    print("="*60)

    # 执行提取
    extractor = DocumentExtractor(force=args.force)
    extracted_data, file_id_mapping = extractor.extract_batch(directory, args.workers)

    # 保存结果
    if extracted_data:
        output_path = Path(args.output)
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(extracted_data, f, ensure_ascii=False, indent=2)

        mapping_path = Path(args.mapping)
        with open(mapping_path, 'w', encoding='utf-8') as f:
            json.dump(file_id_mapping, f, ensure_ascii=False, indent=2)

        print(f"\n✅ 提取完成！")
        print(f"📊 统计信息:")
        print(f"   - 总文件数: {extractor.stats['total']}")
        print(f"   - 成功提取: {extractor.stats['extracted']}")
        print(f"   - 跳过文件: {extractor.stats['skipped']}")
        print(f"   - 失败文件: {extractor.stats['failed']}")
        print(f"\n📁 输出文件:")
        print(f"   - {output_path.absolute()}")
        print(f"   - {mapping_path.absolute()}")

        if extractor.stats['errors']:
            print(f"\n⚠️  失败详情:")
            for err in extractor.stats['errors'][:5]:  # 只显示前5个错误
                print(f"   - {Path(err['file']).name}: {err['error']}")
            if len(extractor.stats['errors']) > 5:
                print(f"   ... 还有 {len(extractor.stats['errors']) - 5} 个错误")

    else:
        print("\n⚠️  没有提取到任何内容")
        return 1

    return 0


if __name__ == '__main__':
    sys.exit(main())
