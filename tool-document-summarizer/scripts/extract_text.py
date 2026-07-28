"""Extract source text with explicit coverage metadata and no source writes."""

from __future__ import annotations

import argparse
import hashlib
import json
import sys
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

SUPPORTED = {".pdf", ".docx", ".pptx", ".xlsx", ".xlsm", ".md", ".txt"}


def source_id(path):
    return hashlib.sha256(str(path.resolve()).encode("utf-8")).hexdigest()[:16]


def extract(path):
    suffix = path.suffix.lower()
    units, limits = [], []
    if suffix in {".md", ".txt"}:
        units = [path.read_text(encoding="utf-8")]
        unit_name = "file"
    elif suffix == ".pdf":
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        units = [(page.extract_text() or "") for page in reader.pages]
        unit_name = "page"
    elif suffix == ".docx":
        from docx import Document
        document = Document(str(path))
        units = [paragraph.text for paragraph in document.paragraphs]
        for table in document.tables:
            units.extend(" | ".join(cell.text for cell in row.cells) for row in table.rows)
        unit_name = "paragraph_or_table_row"
    elif suffix == ".pptx":
        from pptx import Presentation
        presentation = Presentation(str(path))
        units = [
            shape.text
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        ]
        unit_name = "shape"
    elif suffix in {".xlsx", ".xlsm"}:
        from openpyxl import load_workbook
        workbook = load_workbook(str(path), read_only=True, data_only=True)
        for sheet in workbook.worksheets:
            units.append(f"[sheet: {sheet.title}]")
            units.extend(
                " | ".join(str(value) for value in row if value is not None)
                for row in sheet.iter_rows(values_only=True)
            )
        workbook.close()
        unit_name = "sheet_or_row"
    else:
        raise ValueError(f"unsupported file type: {suffix}")
    content = "\n".join(unit for unit in units if unit and unit.strip()).strip()
    if not content:
        raise ValueError("extracted content is empty")
    return {
        "id": source_id(path),
        "filename": str(path.resolve()),
        "content": content,
        "extraction": {
            "unit": unit_name,
            "unit_count": len(units),
            "truncated": False,
            "limits": limits,
        },
    }


def main() -> int:
    parser = argparse.ArgumentParser(description="Extract supported documents")
    parser.add_argument("--dir", required=True, type=Path)
    parser.add_argument("--workers", type=int, default=4)
    parser.add_argument("--output", required=True, type=Path)
    parser.add_argument("--mapping", required=True, type=Path)
    args = parser.parse_args()
    if not args.dir.is_dir() or args.workers < 1:
        print("ERROR[input]: --dir must be a directory and --workers must be positive", file=sys.stderr)
        return 2

    files = sorted(
        path
        for path in args.dir.rglob("*")
        if path.is_file()
        and path.suffix.lower() in SUPPORTED
        and not path.name.startswith((".", "~$"))
    )
    if not files:
        print("ERROR[input]: no supported documents found", file=sys.stderr)
        return 1

    results, errors = [], []
    with ThreadPoolExecutor(max_workers=args.workers) as pool:
        futures = {pool.submit(extract, path): path for path in files}
        for future in as_completed(futures):
            path = futures[future]
            try:
                results.append(future.result())
            except Exception as exc:
                errors.append({"file": str(path), "category": type(exc).__name__, "message": str(exc)})
    if errors:
        print(json.dumps({"status": "blocked", "read_errors": errors}, ensure_ascii=False, indent=2), file=sys.stderr)
        return 1

    results.sort(key=lambda item: item["filename"])
    ids = [item["id"] for item in results]
    if len(ids) != len(set(ids)):
        print("ERROR[schema]: duplicate source id", file=sys.stderr)
        return 1
    mapping = {item["id"]: item["filename"] for item in results}
    try:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.mapping.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(json.dumps(results, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        args.mapping.write_text(json.dumps(mapping, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    except OSError as exc:
        print(f"ERROR[write]: {exc}", file=sys.stderr)
        return 1
    print(json.dumps({"documents": len(results), "output": str(args.output.resolve()), "mapping": str(args.mapping.resolve())}, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
