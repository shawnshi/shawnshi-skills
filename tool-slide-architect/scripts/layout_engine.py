"""
<!-- Standard Header -->
@Input: Slide Blueprint (extracted from outline.md)
@Output: Slide Object (via python-pptx)
@Phase: Phase 4 - Responsive Geometry & Object Rendering
@Maintenance Protocol: Layouts optimized for 16:9 (13.333" x 7.5").
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

class LayoutEngine:
    def __init__(self, prs):
        self.prs = prs
        # Mentat Strategy Palette
        self.NAVY = RGBColor(0, 32, 96)
        self.NAVY_LIGHT = RGBColor(0, 51, 153)
        self.BLACK = RGBColor(0, 0, 0)
        self.GRAY = RGBColor(89, 89, 89)
        self.LIGHT_GRAY = RGBColor(242, 242, 242)
        self.WHITE = RGBColor(255, 255, 255)
        
        # 16:9 Geometry (Inches)
        self.W = 13.333
        self.H = 7.5
        self.MARGIN_X = 0.5
        self.MARGIN_Y = 0.4
        self.CONTENT_W = self.W - (2 * self.MARGIN_X)

    def _apply_navy_gradient(self, slide):
        """物理实现海军蓝线性渐变背景"""
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = self.NAVY
        fill.gradient_stops[1].color.rgb = self.NAVY_LIGHT
        fill.gradient_angle = 90

    def _add_frosted_glass_container(self, slide, x, y, w, h):
        """物理模拟磨砂玻璃容器：半透明白色填充 + 细线描边"""
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = self.WHITE
        # python-pptx transparency is 0.0 to 1.0 (0 is opaque)
        # Note: Some versions use fill.fore_color.alpha or similar, 
        # but the standard way is through XML manipulation for precise alpha.
        # Here we use a very light gray as a fallback for standard API.
        shape.line.color.rgb = self.WHITE
        shape.line.width = Pt(0.5)
        return shape

    def _add_text_box(self, slide, text, x, y, w, h, font_size=Pt(18), bold=False, color=None, alignment=None):
        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True
        if not text:
            return txBox
            
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.bold = bold
        if color:
            p.font.color.rgb = color
        if alignment:
            p.alignment = alignment
        return txBox

    def render_slide(self, slide_data, slide_index=1, total_slides=1):
        # Create a blank slide
        slide_layout = self.prs.slide_layouts[6] # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 1. Background Logic: Gradient for Cover and End slide
        is_special = (slide_index == 1 or slide_index == total_slides)
        text_color = self.WHITE if is_special else self.NAVY
        sub_text_color = self.LIGHT_GRAY if is_special else self.GRAY

        if is_special:
            self._apply_navy_gradient(slide)

        # 2. Render Kicker (Action Title / Headline)
        self._add_text_box(slide, slide_data.get('kicker', ''), 
                          Inches(self.MARGIN_X), Inches(0.3 if not is_special else 2.5), 
                          Inches(12), Inches(0.8 if not is_special else 1.5), 
                          font_size=Pt(32 if not is_special else 44), 
                          bold=True, color=text_color,
                          alignment=PP_ALIGN.CENTER if is_special else None)
        
        # 3. Render Lead-in (Sub-headline)
        if not is_special:
            self._add_text_box(slide, slide_data.get('lead_in', ''), 
                              Inches(self.MARGIN_X), Inches(1.1), Inches(12), Inches(0.4), 
                              font_size=Pt(16), color=sub_text_color)
        else:
            self._add_text_box(slide, slide_data.get('lead_in', ''), 
                              Inches(self.MARGIN_X), Inches(4.0), Inches(12), Inches(0.4), 
                              font_size=Pt(20), color=sub_text_color, alignment=PP_ALIGN.CENTER)

        # 4. Smart Layout Dispatcher (Skip for cover/end)
        if not is_special:
            layout_desc = slide_data.get('layout_desc', '').lower()
            if '左右' in layout_desc or 'split' in layout_desc or '30%' in layout_desc:
                self._render_split_layout(slide, slide_data)
            else:
                self._render_default_layout(slide, slide_data)

            # 5. Render Bumper (The So-What)
            bumper = slide_data.get('bumper', '')
            if bumper:
                # Use frosted glass style for bumper on content slides
                self._add_frosted_glass_container(slide, Inches(self.MARGIN_X), Inches(6.4), 
                                                Inches(self.CONTENT_W), Inches(0.6))
                self._add_text_box(slide, f"➤ {bumper}", 
                                  Inches(self.MARGIN_X + 0.2), Inches(6.45), 
                                  Inches(self.CONTENT_W - 0.4), Inches(0.5), 
                                  font_size=Pt(18), bold=True, color=self.NAVY)

        # 6. Trust Anchor / Footer
        trust = slide_data.get('trust_anchor', '')
        if trust and not is_special:
             self._add_text_box(slide, f"Source: {trust}", 
                               Inches(self.MARGIN_X), Inches(7.1), Inches(8), Inches(0.3), 
                               font_size=Pt(10), color=self.GRAY)

    def _render_default_layout(self, slide, slide_data):
        # Main Body Area
        body_text = slide_data.get('body', '')
        self._add_text_box(slide, body_text, 
                          Inches(self.MARGIN_X), Inches(1.8), Inches(self.CONTENT_W), Inches(4.0), 
                          font_size=Pt(20))

    def _render_split_layout(self, slide, slide_data):
        # Left 40% Text, Right 60% Frosted Glass Visual
        body_text = slide_data.get('body', '')
        # Left Content
        self._add_text_box(slide, body_text, 
                          Inches(self.MARGIN_X), Inches(1.8), Inches(self.W * 0.4), Inches(4.0), 
                          font_size=Pt(18))
        
        # Right Frosted Glass Container
        visual_desc = slide_data.get('visual_desc', 'Visual Evidence')
        self._add_frosted_glass_container(slide, Inches(self.W * 0.45), Inches(1.8), 
                                        Inches(self.W * 0.5), Inches(4.0))
        
        self._add_text_box(slide, f"[Visual Area: {visual_desc}]", 
                          Inches(self.W * 0.45), Inches(3.5), 
                          Inches(self.W * 0.5), Inches(1.0), 
                          font_size=Pt(14), color=self.GRAY, alignment=PP_ALIGN.CENTER)
